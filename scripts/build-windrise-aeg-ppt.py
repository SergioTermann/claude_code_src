#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成 Windrise 创新点 + AEG 演示 PPT(深色工业风)。依赖 python-pptx。
输出: 项目根/Windrise-创新点与AEG.pptx"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Windrise-创新点与AEG.pptx"

# ---- palette ----
BG      = RGBColor(0x12, 0x1C, 0x28)
CARD    = RGBColor(0x1A, 0x27, 0x34)
CARD2   = RGBColor(0x22, 0x31, 0x40)
INK     = RGBColor(0xEA, 0xF0, 0xF5)
INK2    = RGBColor(0xA6, 0xB6, 0xC4)
INK3    = RGBColor(0x74, 0x86, 0x96)
TEAL    = RGBColor(0x35, 0xC6, 0xBA)
TEALDK  = RGBColor(0x12, 0x85, 0x7C)
AMBER   = RGBColor(0xE0, 0xA9, 0x4A)
GREEN   = RGBColor(0x57, 0xC2, 0x8C)
LINE    = RGBColor(0x2A, 0x3A, 0x48)
CJK = "微软雅黑"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_cjk(run, font):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)


def slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=CJK, anchor=MSO_ANCHOR.TOP, spacing=1.0, space_after=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for line_runs in runs if isinstance(runs[0], list) else [runs]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(space_after)
        for seg, opt in line_runs:
            r = p.add_run(); r.text = seg
            r.font.size = Pt(opt.get('size', size))
            r.font.bold = opt.get('bold', bold)
            r.font.color.rgb = opt.get('color', color)
            fn = opt.get('font', font)
            r.font.name = fn; _set_cjk(r, fn)
    return tb


def eyebrow(s, x, y, txt, color=TEAL):
    rect(s, x, y + 0.02, 0.32, 0.03, fill=color)
    text(s, x + 0.42, y - 0.14, 8, 0.35, txt, size=12, color=color, bold=True, font=MONO)


# ============================================================ 1 COVER
s = slide()
rect(s, 0, 0, 13.333/1, 0.14, fill=TEAL)  # top bar
eyebrow(s, 0.9, 1.5, "WIND TURBINE FAULT DIAGNOSIS · 风电运维智能诊断")
text(s, 0.9, 2.0, 11.5, 2.2,
     [[("Windrise 风起时域", {'size': 44, 'bold': True, 'color': INK})],
      [("面向风电机组故障处理的本地化知识构建与检索增强问答系统", {'size': 22, 'color': INK2})]],
     spacing=1.15, space_after=10)
text(s, 0.9, 4.2, 11.4, 1.2,
     "把本地风电资料构建为带机型约束、来源证据与故障机理的领域知识；由本地/内网模型在该知识约束下生成可核查、可续推的现场排查答案。",
     size=17, color=INK2, spacing=1.3)
# stat chips
chips = [("11,865", "故障记录", TEAL), ("2,491", "知识图谱节点", TEAL),
         ("100%", "机理·验证闭环", GREEN), ("本地/离线", "数据不外发", AMBER)]
cx = 0.9
for n, k, c in chips:
    rect(s, cx, 5.6, 2.75, 1.05, fill=CARD, line=LINE, radius=True)
    text(s, cx + 0.25, 5.75, 2.4, 0.5, n, size=26, color=c, bold=True, font=MONO)
    text(s, cx + 0.25, 6.28, 2.4, 0.4, k, size=13, color=INK2)
    cx += 2.95

# ============================================================ 2 PAINS
s = slide()
eyebrow(s, 0.9, 0.7, "01 · 现有方式解决不了的现场难题")
text(s, 0.9, 1.1, 11.5, 0.9, "五类现有方法都只处理文本相似度，未利用运维关系",
     size=28, color=INK, bold=True)
pains = [
    ("同码异义", "同一数字故障码在不同机型下含义不同，单码最多 19 种含义，按数字检索极易查错手册。", AMBER),
    ("资料分散格式不一", "厂家/机型/风场资料格式各异(PDF/Word/RTF/CSV)，难以统一组织。", AMBER),
    ("检索无法表达关系", "答不出'属于哪个系统、哪个部件、由什么原因触发、应先查什么、依据来自哪份资料'。", AMBER),
    ("通用模型答案无依据", "语言流畅但缺本地约束，不可追溯，现场专工无法审核。", AMBER),
    ("单轮难支持多轮排查", "现场靠'正常/异常/压力上不来'短反馈推进，单轮检索无法据此续推。", AMBER),
    ("数据外发风险", "风场资料与检修经验属企业内部，依赖公网模型有外发风险。", AMBER),
]
x0, y0, w, h, gx, gy = 0.9, 2.15, 3.75, 1.9, 0.1, 0.15
for i, (t, d, c) in enumerate(pains):
    col, row = i % 3, i // 3
    x = x0 + col * (w + gx); y = y0 + row * (h + gy)
    rect(s, x, y, w, h, fill=CARD, line=LINE, radius=True)
    rect(s, x, y, 0.06, h, fill=c)
    text(s, x + 0.25, y + 0.2, w - 0.45, 0.5, t, size=16, color=INK, bold=True)
    text(s, x + 0.25, y + 0.72, w - 0.45, 1.05, d, size=12.5, color=INK2, spacing=1.15)

# ============================================================ 3 SIX LAYERS
s = slide()
eyebrow(s, 0.9, 0.7, "02 · 六层核心创新")
text(s, 0.9, 1.1, 11.5, 0.9, "从资料到多轮排查，构成一条完整诊断流水线",
     size=28, color=INK, bold=True)
layers = [
    ("L1", "本地化领域知识构建", "抽取故障码·机型·系统·部件·原因·处理·来源，归一化+别名保留(5645 别名)", TEAL),
    ("L2", "同码异义消歧", "结合风场/机型/上下文限定检索范围，范围不明主动提示并给区分条件", TEAL),
    ("L3", "故障机理归纳", "挂接液压/机械/电气/信号/通信/保护 6 类机理原型，答'为什么、验证哪个量'", TEAL),
    ("L4", "鉴别诊断与反事实验证", "机理竞争时给区分证据与反事实试验，把换件经验升级为可证伪根因", TEAL),
    ("L5", "证据回溯", "结论可回溯到本地条目/机型/部件/来源资料，现场专工可核对", TEAL),
    ("L6", "多轮短反馈排查", "保存状态、短反馈续推、话题归拢，结构化输出'判断→验证→标准→反馈'", TEAL),
]
y = 1.95
for lid, t, d, c in layers:
    rect(s, 0.9, y, 11.5, 0.82, fill=CARD, line=LINE, radius=True)
    rect(s, 1.05, y + 0.16, 0.75, 0.5, fill=None, line=c, line_w=1.5, radius=True)
    text(s, 1.05, y + 0.16, 0.75, 0.5, lid, size=16, color=c, bold=True, font=MONO,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.0, y + 0.11, 3.4, 0.6, t, size=15.5, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.5, y + 0.11, 6.7, 0.6, d, size=12.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    y += 0.9

# ============================================================ 4 MECHANISM GRAPH
s = slide()
eyebrow(s, 0.9, 0.7, "03 · 机理增强知识图谱")
text(s, 0.9, 1.1, 11.5, 0.9, "把'故障码是什么'升级为'为什么发生、如何验证'",
     size=28, color=INK, bold=True)
text(s, 0.9, 2.0, 6.2, 3.6,
     [[("机理推理节点", {'size': 15, 'bold': True, 'color': TEAL})],
      [("mechanism_archetype · failure_mode · propagation_step · observable · verification_test · control_barrier", {'size': 13, 'color': INK2, 'font': MONO})],
      [("", {'size': 6})],
      [("竞争机理判别节点", {'size': 15, 'bold': True, 'color': TEAL})],
      [("diagnostic_hypothesis · discriminating_evidence · counterfactual_test · decision_rule", {'size': 13, 'color': INK2, 'font': MONO})],
      [("", {'size': 6})],
      [("回答现场真正关心的三问", {'size': 15, 'bold': True, 'color': INK})],
      [("为什么发生？为什么查这个量？如何证根因而非凭经验换件？", {'size': 13.5, 'color': INK2})]],
     spacing=1.25, space_after=6)
gm = [("6", "机理原型"), ("1,686", "机理节点"), ("2,601", "机理关系"),
      ("100%", "机理闭环"), ("100%", "假设鉴别"), ("3.0", "平均路径深度")]
gx0, gy0 = 7.5, 2.05
for i, (n, k) in enumerate(gm):
    col, row = i % 3, i // 3
    x = gx0 + col * 1.7; y = gy0 + row * 1.4
    rect(s, x, y, 1.55, 1.2, fill=CARD, line=LINE, radius=True)
    cc = GREEN if "%" in n else TEAL
    text(s, x, y + 0.2, 1.55, 0.55, n, size=23, color=cc, bold=True, font=MONO, align=PP_ALIGN.CENTER)
    text(s, x, y + 0.78, 1.55, 0.35, k, size=11.5, color=INK2, align=PP_ALIGN.CENTER)

# ============================================================ 5 AEG CONCEPT
s = slide()
rect(s, 0, 0, 13.333, 0.1, fill=TEAL)
eyebrow(s, 0.9, 0.7, "04 · 创新亮点 · 歧义熵门控 AEG")
text(s, 0.9, 1.1, 11.5, 1.0, "用信息论在检索前主动预测消歧风险",
     size=30, color=INK, bold=True)
# formula box
rect(s, 0.9, 2.25, 11.5, 1.15, fill=CARD2, line=TEAL, line_w=1.2, radius=True)
text(s, 1.3, 2.42, 10.7, 0.5, "H(c) = − Σ pᵢ · log₂ pᵢ", size=26, color=TEAL, bold=True, font=MONO)
text(s, 1.3, 2.95, 10.7, 0.4, "pᵢ = 故障码 c 映射到第 i 种含义的概率；H=0 单义直接答，H 越大越该在检索前追问机型。",
     size=13.5, color=INK2)
cols = [
    ("被动 → 主动", "现有做法检索后发现歧义才提示；AEG 在检索之前用一个标量预测风险，主动决定是否追问。"),
    ("可量化 · 可设阈值", "阈值 τ 直接对应'追问成本 ↔ 误命中'的帕累托工作点，单一可解释度量便于审核。"),
    ("零训练 · 零 LLM 调用", "熵由知识库含义分布闭式算出，无需训练模型、无需调用大模型，适合本地/离线部署。"),
]
cx = 0.9
for t, d in cols:
    rect(s, cx, 3.75, 3.75, 2.5, fill=CARD, line=LINE, radius=True)
    rect(s, cx + 0.3, 4.0, 0.5, 0.06, fill=TEAL)
    text(s, cx + 0.3, 4.2, 3.15, 0.7, t, size=16.5, color=INK, bold=True, spacing=1.05)
    text(s, cx + 0.3, 5.0, 3.15, 1.15, d, size=13, color=INK2, spacing=1.2)
    cx += 3.95

# ============================================================ 6 AEG EMPIRICAL
s = slide()
eyebrow(s, 0.9, 0.7, "05 · AEG 有效性 · 留一法真实实测")
text(s, 0.9, 1.1, 11.5, 0.9, "11,476 条真实故障记录，实测而非理论期望",
     size=28, color=INK, bold=True)
kpis = [("51.7%", "留一实测基线误命中", AMBER),
        ("0.92", "熵↔误命中 Spearman", TEAL),
        ("0.72", "熵↔增益 相关(独立验证)", TEAL),
        ("→ 0", "门控 62% 查询后误命中", GREEN)]
cx = 0.9
for n, k, c in kpis:
    rect(s, cx, 2.1, 2.75, 1.35, fill=CARD, line=LINE, radius=True)
    text(s, cx + 0.2, 2.32, 2.4, 0.6, n, size=30, color=c, bold=True, font=MONO)
    text(s, cx + 0.2, 2.95, 2.4, 0.45, k, size=12, color=INK2, spacing=1.05)
    cx += 2.95
# three findings
finds = [
    ("熵预测误命中", "实测误命中率随熵单调上升：单义码 0% → 高熵码(H>2) 90.6%。"),
    ("门控帕累托最优", "相同 30% 追问预算下，熵门控残余误命中 24.4% vs 随机 36.2%，多消除 32.5%。"),
    ("增益由熵驱动(独立)", "'追问机型'的实测增益与熵相关 0.72，高熵码收益是低熵码 3 倍；单义码追问收益为 0。"),
]
y = 3.75
for t, d in finds:
    rect(s, 0.9, y, 11.5, 0.86, fill=CARD, line=LINE, radius=True)
    rect(s, 0.9, y, 0.06, 0.86, fill=TEAL)
    text(s, 1.2, y + 0.13, 3.5, 0.6, t, size=15, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.8, y + 0.13, 7.4, 0.6, d, size=12.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE, spacing=1.05)
    y += 0.95
text(s, 0.9, 6.75, 11.5, 0.4,
     "诚实说明：留一实测比闭式期望(28.8%基线)严峻但结论一致；23.8% 查询因资料不足弃答，不计误命中。",
     size=11.5, color=INK3)

# ============================================================ 7 COMPARISON
s = slide()
eyebrow(s, 0.9, 0.7, "06 · 相比流行方法的优势")
text(s, 0.9, 1.1, 11.5, 0.9, "赢在工程适用性与可解释性，而非模型能力",
     size=28, color=INK, bold=True)
rows = [
    ("方法", "需训练", "需LLM", "时机", "面对不确定"),
    ("普通 RAG / 相似度检索", "否", "生成时", "后置", "直接答(会误命中)"),
    ("Self-RAG / 自适应检索", "是", "是", "检索中", "反思后重检索"),
    ("Corrective RAG", "是", "是", "检索后", "触发替代检索"),
    ("语义熵 / 答案不确定性", "否", "多次采样", "后置", "标记幻觉"),
    ("置信度弃答", "常是", "—", "后置", "弃答给人"),
    ("AEG（本方法）", "否", "否", "检索前", "主动问机型消除不确定"),
]
tw, th = 11.5, 4.4
gtbl = s.shapes.add_table(len(rows), 5, Inches(0.9), Inches(2.05), Inches(tw), Inches(th)).table
gtbl.columns[0].width = Inches(3.7)
for w, cwi in zip(range(1, 5), [1.55, 1.75, 1.6, 2.9]):
    gtbl.columns[w].width = Inches(cwi)
for ri, row in enumerate(rows):
    gtbl.rows[ri].height = Inches(th / len(rows))
    for ci, val in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = TEALDK
        elif ri == len(rows) - 1:
            cell.fill.fore_color.rgb = CARD2
        else:
            cell.fill.fore_color.rgb = CARD if ri % 2 else BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(12.5 if ri else 12.5)
        r.font.bold = (ri == 0 or ri == len(rows) - 1)
        is_aeg = ri == len(rows) - 1
        r.font.color.rgb = INK if ri == 0 else (TEAL if is_aeg else INK2)
        r.font.name = CJK; _set_cjk(r, CJK)
text(s, 0.9, 6.65, 11.5, 0.5,
     "定位：AEG 是更便宜、更透明、更早介入的消歧决策层——零成本信息论标量，在检索前量化'该不该打扰用户'。",
     size=12.5, color=INK2, bold=True)

# ============================================================ 8 VALUE / CLOSE
s = slide()
rect(s, 0, 0, 13.333, 0.14, fill=TEAL)
eyebrow(s, 0.9, 1.2, "07 · 部署与价值")
text(s, 0.9, 1.65, 11.5, 1.0, "能安全上线的本地化风电诊断底座", size=32, color=INK, bold=True)
vals = [
    ("数据不出内网", "本地/内网模型 + 离线嵌入(hash/本地)，风场资料与检修经验不外发。"),
    ("可核查可审计", "答案回到本地资料来源，满足运维审核；机理链支撑可证伪根因。"),
    ("贴合现场节奏", "短反馈续推、每轮只推一个动作，契合登机作业交互。"),
    ("论文/专利就绪", "AEG 有实测支撑、论文章节与专利权利要求齐备。"),
]
x0, y0, w, h = 0.9, 3.0, 5.65, 1.5
for i, (t, d) in enumerate(vals):
    col, row = i % 2, i // 2
    x = x0 + col * (w + 0.2); y = y0 + row * (h + 0.2)
    rect(s, x, y, w, h, fill=CARD, line=LINE, radius=True)
    rect(s, x, y, 0.06, h, fill=TEAL)
    text(s, x + 0.3, y + 0.22, w - 0.5, 0.5, t, size=17, color=INK, bold=True)
    text(s, x + 0.3, y + 0.75, w - 0.5, 0.65, d, size=13, color=INK2, spacing=1.15)
text(s, 0.9, 6.7, 11.5, 0.5, "Windrise 风起时域 · 面向风电机组故障处理的本地化知识构建、检索增强问答系统及方法",
     size=12, color=INK3)

prs.save(str(OUT))
print(f"✅ 生成 {OUT}  共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页")
